import io
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION


def create_pptx(result: dict, inputs: dict) -> io.BytesIO:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    period = inputs.get('period', {})
    cc = inputs.get('country_config', {})
    country = cc.get('country', 'Country')
    area = cc.get('area', 'Area')
    currency = cc.get('currency', 'LCU')
    baseline = period.get('baseline_year', 2025)
    t1 = period.get('target1_year', 2030)
    t2 = period.get('target2_year', 2040)
    years = result['years']

    # Filter to forecast years only
    forecast_start = baseline
    fi = [i for i, y in enumerate(years) if y >= forecast_start]
    fy = [years[i] for i in fi]

    # Colors
    BLUE = RGBColor(0x00, 0x22, 0x44)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    GRAY = RGBColor(0x64, 0x74, 0x8B)
    RED = RGBColor(0xEF, 0x44, 0x44)
    GREEN = RGBColor(0x10, 0xB9, 0x81)

    def add_title_slide(title, subtitle):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = BLUE
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
        return slide

    def add_section_slide(title):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(0xF1, 0xF5, 0xF9)
        txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.color.rgb = BLUE
        p.font.bold = True
        return slide

    def add_table_slide(title, headers, rows):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # Title
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.color.rgb = BLUE
        p.font.bold = True

        n_rows = len(rows) + 1
        n_cols = len(headers)
        tbl_width = Inches(12)
        tbl_height = Inches(0.4 * n_rows)
        table = slide.shapes.add_table(n_rows, n_cols, Inches(0.5), Inches(1.2), tbl_width, tbl_height).table

        # Header row
        for j, h in enumerate(headers):
            cell = table.cell(0, j)
            cell.text = str(h)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.bold = True
                paragraph.font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLUE

        # Data rows
        for i, row in enumerate(rows):
            for j, val in enumerate(row):
                cell = table.cell(i + 1, j)
                cell.text = str(val)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(10)
                    paragraph.alignment = PP_ALIGN.RIGHT if j > 0 else PP_ALIGN.LEFT
                if i % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)

        return slide

    def fmt(val, dec=2, div=1):
        return f"{val / div:,.{dec}f}"

    # === SLIDE 1: Title ===
    add_title_slide(
        "WSS Strategic Scenarios Analysis",
        f"{country} — {area} | Baseline {baseline} | Targets {t1} & {t2}"
    )

    # Summary years
    sy = [baseline, t1, t2]
    si = [years.index(y) for y in sy if y in years]

    for sector_key, sector_name in [('water_supply', 'Water Supply'), ('sanitation', 'Sanitation')]:
        sec = result[sector_key]

        # === Section slide ===
        add_section_slide(f"{sector_name}")

        # === Summary table ===
        headers = ['Metric'] + [str(years[i]) for i in si]
        rows = [
            ['Total HH (millions)'] + [fmt(result['total_hh'][i], 3) for i in si],
            ['Target Safely Managed HH (mill)'] + [fmt(sec['target_hh_serv'][0][i], 4) for i in si],
            ['BAU Safely Managed HH (mill)'] + [fmt(sec['bau_hh_serv'][0][i], 4) for i in si],
            ['Service Gap (mill HH)'] + [fmt(sec['service_gap'][i], 4) for i in si],
            ['Investment Need ({currency} bill)'] + [fmt(sec['investment_need'][i], 2, 1000) for i in si],
            ['BAU Investment ({currency} bill)'] + [fmt(sec['bau_investment'][i], 2, 1000) for i in si],
            ['Financing Gap ({currency} bill)'] + [fmt(sec['financing_gap'][i], 2, 1000) for i in si],
        ]
        add_table_slide(f"{sector_name} — Summary", headers, rows)

        # === Intervention impact table ===
        interv = sec['interventions']
        interv_names = list(interv.keys())
        int_headers = ['Intervention'] + [str(years[i]) for i in si]
        int_rows = []
        for name in interv_names:
            display_name = name.replace('_', ' ').title()
            cum_hh = interv[name].get('cumulative_hh', [0] * len(years))
            int_rows.append([display_name] + [fmt(cum_hh[i], 4) for i in si])
        add_table_slide(f"{sector_name} — Intervention Impact (Cumulative HH, millions)", int_headers, int_rows)

    # Save
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output
